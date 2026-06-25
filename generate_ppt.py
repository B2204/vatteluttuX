from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── COLORS ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE  = RGBColor(0x1F, 0x39, 0x64)   # deep navy  – headings
MID_BLUE   = RGBColor(0x2E, 0x74, 0xB5)   # accent blue – sub-headings / rules
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF7)   # table header background
GOLD       = RGBColor(0xC0, 0x8C, 0x00)   # highlight accent
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)   # alternate table rows

FONT = "Times New Roman"

# ─── PRESENTATION SETUP ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

# ─── HELPER UTILITIES ─────────────────────────────────────────────────────────

def white_bg(slide):
    """Fill slide background white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, italic=False,
                 color=BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = FONT
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txBox


def add_paragraph(tf, text, font_size=14, bold=False, italic=False,
                  color=BLACK, align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.alignment  = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text           = text
    run.font.name      = FONT
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return p


def slide_header(slide, title, subtitle=None):
    """Draws a top banner with title (and optional subtitle)."""
    # banner rectangle
    banner = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), prs.slide_width, Inches(1.25)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_BLUE
    banner.line.fill.background()

    add_text_box(slide, title,
                 Inches(0.3), Inches(0.08), Inches(12.5), Inches(0.75),
                 font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.3), Inches(0.78), Inches(12.5), Inches(0.4),
                     font_size=16, italic=True, color=RGBColor(0xAD, 0xD8, 0xE6),
                     align=PP_ALIGN.LEFT)


def add_table(slide, headers, rows, left, top, width, height,
              header_bg=LIGHT_BLUE, alt_row=True):
    """Add a styled table."""
    cols  = len(headers)
    nrows = len(rows) + 1
    tbl   = slide.shapes.add_table(nrows, cols, left, top, width, height).table

    col_w = width // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    # header row
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        tf   = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run  = tf.paragraphs[0].add_run()
        run.text           = h
        run.font.name      = FONT
        run.font.size      = Pt(13)
        run.font.bold      = True
        run.font.color.rgb = WHITE

    # data rows
    for ri, row in enumerate(rows):
        bg = LIGHT_GREY if (alt_row and ri % 2 == 0) else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf   = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run  = tf.paragraphs[0].add_run()
            run.text           = str(val)
            run.font.name      = FONT
            run.font.size      = Pt(12)
            run.font.color.rgb = BLACK
    return tbl


def bullet_box(slide, items, left, top, width, height,
               font_size=15, title=None, title_color=MID_BLUE, title_size=17):
    """A text box with an optional bold title and bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text           = title
        run.font.name      = FONT
        run.font.size      = Pt(title_size)
        run.font.bold      = True
        run.font.color.rgb = title_color
        first = False
    else:
        first = True

    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment   = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        p.level        = 0
        run = p.add_run()
        run.text           = item
        run.font.name      = FONT
        run.font.size      = Pt(font_size)
        run.font.color.rgb = BLACK
    return txBox


def divider(slide, top_inches, color=MID_BLUE):
    line = slide.shapes.add_shape(
        1, Inches(0.3), Inches(top_inches),
        Inches(12.73), Pt(2)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)

# Big center block
title_box = slide.shapes.add_shape(
    1, Inches(0), Inches(1.8), prs.slide_width, Inches(2.8)
)
title_box.fill.solid()
title_box.fill.fore_color.rgb = DARK_BLUE
title_box.line.fill.background()

add_text_box(slide,
    "VattalettuX",
    Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.0),
    font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "A Deep Learning System for Reading Vatteluttu – An Ancient Tamil Script",
    Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.75),
    font_size=20, italic=True, color=RGBColor(0xAD, 0xD8, 0xE6),
    align=PP_ALIGN.CENTER)

# Thin gold line
sep = slide.shapes.add_shape(
    1, Inches(3.5), Inches(4.75), Inches(6.3), Pt(3)
)
sep.fill.solid()
sep.fill.fore_color.rgb = GOLD
sep.line.fill.background()

add_text_box(slide,
    "A Deep Learning-Based OCR System for Ancient Tamil Epigraphy",
    Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.45),
    font_size=16, color=DARK_BLUE, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "[Your Name]  |  PG Department of Computer Applications",
    Inches(0.5), Inches(5.42), Inches(12.3), Inches(0.35),
    font_size=14, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)

add_text_box(slide,
    "[Your College Name], Tamil Nadu, India  |  2025",
    Inches(0.5), Inches(5.78), Inches(12.3), Inches(0.35),
    font_size=14, color=RGBColor(0x44, 0x44, 0x44), align=PP_ALIGN.CENTER)

# footer tags
add_text_box(slide,
    "Vatteluttu  •  Ancient Script  •  OCR  •  CNN  •  Deep Learning  •  Tamil Epigraphy",
    Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
    font_size=12, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – INTRODUCTION 1
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Introduction", "What is Vatteluttu?")

bullet_box(slide, [
    "▸  Vatteluttu is one of the oldest Tamil writing systems",
    "▸  Used from 3rd to 12th century CE by Chola, Pandya & Pallava kingdoms",
    "▸  Thousands of stone inscriptions survive on temple walls, pillars, and rock faces",
    "▸  These inscriptions hold records of royal orders, land grants & daily life of ancient India",
    "▸  Today, only a very small number of trained epigraphists can read this script",
    "▸  Physical erosion of stone surfaces over centuries makes reading even harder",
    "▸  Government agencies (e.g., ASI) have recorded inscriptions, but\n"
    "    converting them to digital text remains fully manual and expert-dependent",
],
    left=Inches(0.4), top=Inches(1.38), width=Inches(7.5), height=Inches(5.5),
    font_size=15)

# Problem statement box (right side)
prob = slide.shapes.add_shape(
    1, Inches(8.2), Inches(1.6), Inches(4.8), Inches(4.5)
)
prob.fill.solid()
prob.fill.fore_color.rgb = RGBColor(0xEB, 0xF3, 0xFB)
prob.line.color.rgb      = MID_BLUE

add_text_box(slide, "🔴 The Problem",
             Inches(8.4), Inches(1.75), Inches(4.4), Inches(0.45),
             font_size=16, bold=True, color=DARK_BLUE)

add_text_box(slide,
    "Historical knowledge stored inside Vatteluttu inscriptions "
    "remains inaccessible to most researchers and the general public "
    "due to a lack of readable form and scarcity of experts.",
    Inches(8.4), Inches(2.25), Inches(4.4), Inches(2.2),
    font_size=14, italic=True, color=BLACK)

add_text_box(slide, "➜  VattalettuX solves this with Artificial Intelligence.",
             Inches(8.4), Inches(4.5), Inches(4.4), Inches(0.6),
             font_size=13, bold=True, color=MID_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – INTRODUCTION 2
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Introduction", "Why Do We Need VattalettuX?")

# Three-problem table
add_text_box(slide, "Three Core Challenges:",
             Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.4),
             font_size=16, bold=True, color=DARK_BLUE)

add_table(slide,
    ["#", "Challenge", "Explanation"],
    [
        ["1", "Noisy Images",
         "Stone photos have cracks, shadows, moss, uneven lighting & erosion"],
        ["2", "Large Character Set",
         "Vatteluttu has 247 distinct character classes — far more than any existing study"],
        ["3", "Output Usability",
         "Results must be in Modern Tamil Unicode so historians can actually use them"],
    ],
    left=Inches(0.4), top=Inches(1.9), width=Inches(12.5), height=Inches(1.6))

divider(slide, 3.65)

add_text_box(slide, "Our Solution — VattalettuX:",
             Inches(0.4), Inches(3.75), Inches(12.5), Inches(0.4),
             font_size=16, bold=True, color=DARK_BLUE)

bullet_box(slide, [
    "▸  A fully web-based OCR system powered by Deep Learning (ResNet CNN)",
    "▸  Automatically reads Vatteluttu characters from stone inscription photos",
    "▸  Converts them to Modern Tamil Unicode text instantly",
    "▸  No installation needed — works in any web browser",
    "▸  Built with React.js (Frontend)  +  FastAPI (Backend)  +  PyTorch CNN Model",
],
    left=Inches(0.4), top=Inches(4.2), width=Inches(12.5), height=Inches(2.8),
    font_size=15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – LITERATURE REVIEW 1
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Literature Review", "Previous Research on Tamil & Ancient Script Recognition")

add_table(slide,
    ["Author(s)", "Work", "Key Contribution"],
    [
        ["Murugan & Visalakshi [1]\n(Heritage Science, 2024)",
         "Ancient Tamil inscription recognition",
         "Used DRL framework — showed importance of separating detection & recognition stages"],
        ["Gayathri Devi et al. [2]\n(Comput. Intelligence, 2022)",
         "Tamil palm leaf manuscript OCR",
         "CNN + morphological preprocessing + CCA — similar challenges to worn stone surfaces"],
        ["Vijaya Arjunan et al. [3]\n(JISIS, 2025)",
         "Vatteluttu script recognition (Deep Learning)",
         "98% accuracy on only 28 character classes — directly related but very limited scope"],
    ],
    left=Inches(0.4), top=Inches(1.4), width=Inches(12.5), height=Inches(2.8))

divider(slide, 4.4)

# Research gap highlight
gap = slide.shapes.add_shape(
    1, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.55)
)
gap.fill.solid()
gap.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
gap.line.color.rgb      = GOLD

add_text_box(slide, "⚠  Research Gap — What Was Missing:",
             Inches(0.6), Inches(4.6), Inches(12.1), Inches(0.4),
             font_size=16, bold=True, color=RGBColor(0x7B, 0x60, 0x00))

bullet_box(slide, [
    "▸  Best existing Vatteluttu study covers only 28 character classes — far from the complete script",
    "▸  No study has attempted all 247 Vatteluttu characters",
    "▸  No fully deployed web application exists for real-world use",
    "▸  VattalettuX closes this gap — 247 classes (9× more) + live web application",
],
    left=Inches(0.6), top=Inches(5.05), width=Inches(12.1), height=Inches(1.85),
    font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – LITERATURE REVIEW 2
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Literature Review", "Key Deep Learning Architectures & Techniques Used")

add_table(slide,
    ["Architecture / Technique", "Authors", "Role in VattalettuX"],
    [
        ["LeNet [7]",            "LeCun et al. (1998)",           "Foundation of CNN-based document recognition"],
        ["AlexNet [6]",          "Krizhevsky et al. (2012)",       "Proved deep CNNs outperform traditional methods"],
        ["VGGNet [5]",           "Simonyan & Zisserman (2015)",    "Small 3×3 filters for very deep feature extraction"],
        ["ResNet [4] ★",         "He et al. (2016)",               "Residual/skip connections — our model is built on this"],
        ["Dropout [10]",         "Srivastava et al. (2014)",       "Prevents overfitting — applied at p=0.5 in FC layer"],
        ["Batch Normalization [12]","Ioffe & Szegedy (2015)",      "Faster training & stable convergence"],
        ["Adam Optimizer [11]",  "Kingma & Ba (2015)",             "Adaptive learning rate — used for all training"],
        ["Data Augmentation [16]","Shorten & Khoshgoftaar (2019)", "Strategy for building 2,47,000-image synthetic dataset"],
        ["OpenCV [17]",          "Bradski (2000)",                  "Core image processing — thresholding, CCA, morphology"],
        ["PyTorch [15]",         "Paszke et al. (2019)",           "Deep learning framework for training & inference"],
    ],
    left=Inches(0.4), top=Inches(1.38), width=Inches(12.5), height=Inches(5.6))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – METHODOLOGY 1: SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Methodology", "System Architecture — End-to-End Pipeline")

# Pipeline steps (left column)
steps = [
    ("STEP 1", "User uploads inscription photo via Browser"),
    ("STEP 2", "React.js Frontend sends image to backend via REST API"),
    ("STEP 3", "FastAPI Backend (Python) receives the request"),
    ("STEP 4", "Image Preprocessing  ←  OpenCV (Otsu, Morphology)"),
    ("STEP 5", "Character Segmentation  ←  Connected Component Analysis"),
    ("STEP 6", "CNN Classification  ←  ResNet + PyTorch (247 classes)"),
    ("STEP 7", "Tamil Character Mapping  ←  JSON database lookup"),
    ("STEP 8", "Modern Tamil Unicode output displayed to user"),
]

y = 1.42
for label, desc in steps:
    box = slide.shapes.add_shape(
        1, Inches(0.35), Inches(y), Inches(8.5), Inches(0.6)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE if label not in ("STEP 4","STEP 5","STEP 6","STEP 7") else RGBColor(0xD0, 0xE8, 0xFF)
    box.line.color.rgb = MID_BLUE

    add_text_box(slide, f"  {label}  |  {desc}",
                 Inches(0.4), Inches(y + 0.08), Inches(8.4), Inches(0.48),
                 font_size=13, bold=(label in ("STEP 4","STEP 5","STEP 6","STEP 7")),
                 color=DARK_BLUE)
    y += 0.66

# Right: API endpoints
api = slide.shapes.add_shape(
    1, Inches(9.1), Inches(1.4), Inches(3.9), Inches(4.0)
)
api.fill.solid()
api.fill.fore_color.rgb = RGBColor(0xEB, 0xF3, 0xFB)
api.line.color.rgb = MID_BLUE

add_text_box(slide, "REST API Endpoints",
             Inches(9.2), Inches(1.5), Inches(3.7), Inches(0.4),
             font_size=15, bold=True, color=DARK_BLUE)

for method, ep, desc in [
    ("POST", "/recognize", "Runs full pipeline → returns Tamil text"),
    ("GET",  "/health",    "Checks server & model status"),
    ("GET",  "/character-map", "Returns 247-character mapping database"),
]:
    add_text_box(slide, f"{method}  {ep}",
                 Inches(9.2), Inches(y := y), Inches(3.7), Inches(0.3),
                 font_size=13, bold=True, color=MID_BLUE)
    add_text_box(slide, f"   {desc}",
                 Inches(9.2), Inches(y + 0.3), Inches(3.7), Inches(0.35),
                 font_size=12, color=BLACK)
    y += 0.75

add_text_box(slide,
    "★ Insert: System Architecture Diagram (Fig. 1) here",
    Inches(9.1), Inches(5.65), Inches(3.9), Inches(0.5),
    font_size=11, italic=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – METHODOLOGY 2: PREPROCESSING & SEGMENTATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Methodology", "Stage 1 & 2 — Image Preprocessing & Character Segmentation")

add_text_box(slide, "Stage 1 — Image Preprocessing (4 Steps using OpenCV)",
             Inches(0.4), Inches(1.38), Inches(12.5), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(slide,
    ["Step", "Operation", "Purpose"],
    [
        ["1", "Grayscale Conversion",       "Removes color info — reduces complexity, focuses on shape"],
        ["2", "Otsu's Binarization",         "Converts to black & white — auto-selects best threshold"],
        ["3", "Morphological Opening",       "Erosion → Dilation — removes noise specks from stone surface"],
        ["4", "Histogram Equalization",      "Spreads brightness range — improves contrast for faded characters"],
    ],
    left=Inches(0.4), top=Inches(1.82), width=Inches(12.5), height=Inches(1.7))

add_text_box(slide, "★ Insert: Preprocessing Pipeline Image (Fig. 2) here — Original → Grayscale → Binary → Clean",
             Inches(0.4), Inches(3.6), Inches(12.5), Inches(0.4),
             font_size=12, italic=True, color=GOLD)

divider(slide, 4.1)

add_text_box(slide, "Stage 2 — Character Segmentation (Connected Component Analysis)",
             Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

bullet_box(slide, [
    "▸  Algorithm (Suzuki & Abe [9]) scans binary image and groups touching foreground pixels into blobs",
    "▸  Each blob gets a unique label and a bounding box  [x, y, width, height]",
    "▸  Invalid blobs are filtered out using two tests:",
    "      •  Area Filter:  A_min  ≤  Area(blob)  ≤  A_max          (removes dust specks & giant merges)",
    "      •  Aspect Ratio Filter:  0.1  ≤  Width/Height  ≤  10.0   (removes line fragments)",
    "▸  Valid blobs are cropped from the image and resized to 64×64 pixels for the CNN",
],
    left=Inches(0.4), top=Inches(4.65), width=Inches(12.5), height=Inches(2.6),
    font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – METHODOLOGY 3: CNN CLASSIFIER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Methodology", "Stage 3 — CNN Classifier (ResNet-Based Architecture)")

add_text_box(slide, "Model Architecture:",
             Inches(0.4), Inches(1.38), Inches(7.0), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(slide,
    ["Layer", "Details"],
    [
        ["Input",                   "64 × 64 grayscale character image"],
        ["Residual Blocks 1 – 4",   "Convolutional layers with skip connections + ReLU activation"],
        ["Batch Normalization",     "Applied after every convolutional layer [12]"],
        ["Global Average Pooling",  "Compresses feature maps into a single feature vector"],
        ["Fully Connected Layer",   "512 units"],
        ["Dropout  (p = 0.5)",      "Regularization — prevents overfitting [10]"],
        ["Output Layer",            "247 units with Softmax — one score per Vatteluttu character"],
    ],
    left=Inches(0.4), top=Inches(1.82), width=Inches(7.0), height=Inches(3.4))

add_text_box(slide, "Training Setup:",
             Inches(0.4), Inches(5.35), Inches(7.0), Inches(0.38),
             font_size=14, bold=True, color=DARK_BLUE)
bullet_box(slide, [
    "▸  Framework: PyTorch  |  Optimizer: Adam  (lr = 0.001, β₁=0.9, β₂=0.999)",
    "▸  Batch Size: 32  |  Epochs: 100  |  Scheduler: ReduceLROnPlateau (patience=10)",
    "▸  Augmentation: Rotation ±15°, Brightness jitter, Flipping, Gaussian noise",
],
    left=Inches(0.4), top=Inches(5.78), width=Inches(7.0), height=Inches(1.5),
    font_size=13)

# Right side — CNN diagram placeholder
right = slide.shapes.add_shape(
    1, Inches(7.7), Inches(1.38), Inches(5.3), Inches(5.9)
)
right.fill.solid()
right.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFF)
right.line.color.rgb      = MID_BLUE

add_text_box(slide, "CNN Architecture Diagram",
             Inches(7.9), Inches(1.55), Inches(4.9), Inches(0.45),
             font_size=14, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

layers = [
    ("Input", "64×64 Grayscale", LIGHT_BLUE),
    ("ResBlock 1", "Conv + BN + ReLU + Skip", RGBColor(0xD0,0xE8,0xFF)),
    ("ResBlock 2", "Conv + BN + ReLU + Skip", RGBColor(0xD0,0xE8,0xFF)),
    ("ResBlock 3", "Conv + BN + ReLU + Skip", RGBColor(0xD0,0xE8,0xFF)),
    ("ResBlock 4", "Conv + BN + ReLU + Skip", RGBColor(0xD0,0xE8,0xFF)),
    ("GAP",        "Global Avg Pooling",       RGBColor(0xC8,0xFF,0xC8)),
    ("FC + Dropout","512 units, p=0.5",        RGBColor(0xFF,0xF0,0xC8)),
    ("Output",     "247-class Softmax",        RGBColor(0xFF,0xD0,0xD0)),
]
y_l = 2.1
for name, detail, clr in layers:
    lb = slide.shapes.add_shape(1, Inches(8.0), Inches(y_l), Inches(4.7), Inches(0.55))
    lb.fill.solid()
    lb.fill.fore_color.rgb = clr
    lb.line.color.rgb      = MID_BLUE
    add_text_box(slide, f"{name}  —  {detail}",
                 Inches(8.05), Inches(y_l+0.07), Inches(4.6), Inches(0.42),
                 font_size=12, bold=(name in ("Input","Output")), color=DARK_BLUE,
                 align=PP_ALIGN.CENTER)
    y_l += 0.6
    if y_l < 6.8 and name != "Output":
        arr = slide.shapes.add_shape(1, Inches(10.2), Inches(y_l-0.07),
                                      Inches(0.2), Inches(0.15))
        arr.fill.solid()
        arr.fill.fore_color.rgb = DARK_BLUE
        arr.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – METHODOLOGY 4: DATASET & TAMIL MAPPING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Methodology", "Stage 4 — Synthetic Dataset & Tamil Character Mapping")

# Dataset section
add_text_box(slide, "Building the Training Dataset (No Real Dataset Existed):",
             Inches(0.4), Inches(1.38), Inches(12.5), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

bullet_box(slide, [
    "▸  There is NO publicly available labeled dataset of Vatteluttu characters",
    "▸  Solution: Generated 2,47,000 synthetic images from authentic Vatteluttu fonts",
    "▸  1,000 variation images per character class  ×  247 classes",
    "▸  Augmentation: Random rotation & scaling, Dilation & erosion (stone wear simulation),",
    "    Gaussian + salt-and-pepper noise, Stone-like textured backgrounds",
],
    left=Inches(0.4), top=Inches(1.82), width=Inches(7.5), height=Inches(2.3),
    font_size=14)

add_table(slide,
    ["Split", "Percentage", "Number of Images"],
    [
        ["Training",   "70%", "1,72,900"],
        ["Validation", "15%", "37,050"],
        ["Testing",    "15%", "37,050"],
        ["Total",      "100%", "2,47,000"],
    ],
    left=Inches(0.4), top=Inches(4.2), width=Inches(5.0), height=Inches(2.0))

# Right: character categories table
add_text_box(slide, "247 Character Categories:",
             Inches(5.8), Inches(1.38), Inches(7.0), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(slide,
    ["Category", "Tamil Name", "Classes"],
    [
        ["Vowels",              "உயிர் எழுத்து",  "12"],
        ["Aytham",             "ஆய்தம்",          "1"],
        ["Pure Consonants",    "மெய் எழுத்து",    "18"],
        ["Consonants (with \'a\')", "—",           "18"],
        ["Compound (Uyirmei)", "உயிர்மெய்",       "198"],
        ["TOTAL",              "",                 "247"],
    ],
    left=Inches(5.8), top=Inches(1.82), width=Inches(7.1), height=Inches(2.8))

add_text_box(slide,
    "Tamil Mapping:  CNN output label (e.g. va_037) → JSON lookup → Modern Tamil Unicode character",
    Inches(5.8), Inches(4.7), Inches(7.1), Inches(0.5),
    font_size=13, italic=True, color=DARK_BLUE)

add_text_box(slide, "★ Insert: Dataset pie chart (70/15/15) or sample augmented images here",
             Inches(5.8), Inches(5.3), Inches(7.1), Inches(0.4),
             font_size=11, italic=True, color=GOLD)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – ALGORITHM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Algorithm", "VattalettuX — Step-by-Step Recognition Algorithm")

algo_steps = [
    ("INPUT",    "Inscription image (photograph of a stone carving)",
     RGBColor(0xD0,0xE8,0xFF)),

    ("STEP 1\nPreprocessing",
     "Convert to Grayscale  →  Apply Otsu's Binarization  →  Morphological Opening (noise removal)  →  Histogram Equalization",
     LIGHT_BLUE),

    ("STEP 2\nSegmentation",
     "Run Connected Component Analysis  →  Label all blobs  →  Draw bounding boxes\n"
     "Filter blobs: Area filter  (A_min ≤ Area ≤ A_max)  +  Aspect ratio filter  (0.1 ≤ W/H ≤ 10.0)\n"
     "Crop valid blobs  →  Resize each to 64×64 pixels",
     RGBColor(0xE8,0xF5,0xFF)),

    ("STEP 3\nClassification",
     "Feed 64×64 image into ResNet CNN  →  Get probability scores for all 247 classes\n"
     "Select class with highest Softmax score  →  Record confidence score",
     RGBColor(0xFD,0xF5,0xE0)),

    ("STEP 4\nMapping",
     "Look up CNN output label in JSON character map  →  Retrieve Modern Tamil Unicode equivalent",
     RGBColor(0xE8,0xFF,0xE8)),

    ("OUTPUT",   "Sequence of Modern Tamil Unicode characters + Confidence scores per character",
     RGBColor(0xFF,0xD0,0xD0)),
]

y = 1.38
for label, content, clr in algo_steps:
    h = 0.7 if label in ("INPUT", "OUTPUT") else 0.9
    box = slide.shapes.add_shape(1, Inches(0.35), Inches(y), Inches(12.6), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = clr
    box.line.color.rgb      = MID_BLUE

    add_text_box(slide, label,
                 Inches(0.45), Inches(y + 0.07), Inches(1.6), Inches(h - 0.1),
                 font_size=12, bold=True, color=DARK_BLUE)

    add_text_box(slide, content,
                 Inches(2.1), Inches(y + 0.05), Inches(10.7), Inches(h - 0.1),
                 font_size=13, color=BLACK)
    y += h + 0.1


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – IMPLEMENTATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Implementation", "Technology Stack & Web Application Features")

add_text_box(slide, "Technology Stack:",
             Inches(0.4), Inches(1.38), Inches(6.2), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(slide,
    ["Component", "Technology"],
    [
        ["Frontend UI",       "React.js + TypeScript"],
        ["Backend Server",    "FastAPI (Python)"],
        ["Deep Learning",     "PyTorch"],
        ["Image Processing",  "OpenCV"],
        ["Model Architecture","ResNet-inspired CNN"],
        ["Character Database","JSON  (247 entries)"],
        ["Deployment",        "Web Application (any browser, no install)"],
    ],
    left=Inches(0.4), top=Inches(1.82), width=Inches(6.2), height=Inches(3.5))

add_text_box(slide, "Web Application Features:",
             Inches(0.4), Inches(5.45), Inches(6.2), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

bullet_box(slide, [
    "▸  Drag-and-drop image upload   ▸  Live bounding box preview",
    "▸  Character-by-character result display   ▸  Confidence scores shown",
    "▸  Export output as plain text file",
],
    left=Inches(0.4), top=Inches(5.88), width=Inches(6.2), height=Inches(1.3),
    font_size=13)

# Right — screenshot placeholder
ph = slide.shapes.add_shape(
    1, Inches(6.9), Inches(1.38), Inches(6.1), Inches(5.9)
)
ph.fill.solid()
ph.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
ph.line.color.rgb      = MID_BLUE

add_text_box(slide, "★ Insert: Web Application Screenshot here",
             Inches(7.1), Inches(3.5), Inches(5.7), Inches(0.5),
             font_size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "Upload Panel  |  Bounding Box Preview\n"
    "Recognized Characters  |  Tamil Unicode Output\n"
    "Confidence Scores  |  Export Button",
    Inches(7.1), Inches(4.1), Inches(5.7), Inches(1.5),
    font_size=12, italic=True, color=RGBColor(0x88,0x88,0x88),
    align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – RESULTS 1: CLASSIFICATION ACCURACY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Result and Analysis", "Classification Accuracy — 37,050 Held-Out Test Images")

# Headline stats
for val, lbl, x in [
    ("92.8%", "Overall Top-1 Accuracy", 1.2),
    ("98.1%", "Overall Top-5 Accuracy", 5.1),
    ("247",   "Character Classes",       9.0),
]:
    box = slide.shapes.add_shape(1, Inches(x), Inches(1.38), Inches(3.0), Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = DARK_BLUE
    box.line.fill.background()
    add_text_box(slide, val,
                 Inches(x), Inches(1.42), Inches(3.0), Inches(0.7),
                 font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, lbl,
                 Inches(x), Inches(2.1), Inches(3.0), Inches(0.38),
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_table(slide,
    ["Character Category", "No. of Classes", "Top-1 Accuracy", "Top-5 Accuracy"],
    [
        ["Vowels  (உயிர்)",             "12",  "97.3%",  "99.8%"],
        ["Aytham  (ஆய்தம்)",           "1",   "99.1%",  "100.0%"],
        ["Pure Consonants  (மெய்)",     "18",  "95.6%",  "99.2%"],
        ["Consonants  (with \'a\')",    "18",  "96.1%",  "99.4%"],
        ["Compound Characters (உயிர்மெய்)", "198", "91.4%", "97.8%"],
        ["OVERALL",                     "247", "92.8%",  "98.1%"],
    ],
    left=Inches(0.4), top=Inches(2.65), width=Inches(8.5), height=Inches(3.6))

add_text_box(slide, "★ Insert: Accuracy bar chart (by category) here",
             Inches(9.1), Inches(2.65), Inches(3.9), Inches(3.6),
             font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "Note: Compound (Uyirmei) characters are harder — 198 classes differing only in small diacritic vowel marks",
    Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.4),
    font_size=12, italic=True, color=RGBColor(0x55,0x55,0x55))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – RESULTS 2: SEGMENTATION & COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Result and Analysis", "Segmentation Results & Comparison with Related Work")

add_text_box(slide, "Segmentation Performance (CCA):",
             Inches(0.4), Inches(1.38), Inches(6.5), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(slide,
    ["Image Condition", "Precision", "Recall", "F1-Score"],
    [
        ["Well-separated characters",  "0.94", "0.96", "0.95"],
        ["Moderately spaced",          "0.88", "0.91", "0.89"],
        ["Closely packed characters",  "0.79", "0.83", "0.81"],
        ["Overall Average",            "0.87", "0.90", "0.88"],
    ],
    left=Inches(0.4), top=Inches(1.82), width=Inches(6.5), height=Inches(2.0))

add_text_box(slide, "★ Insert: Segmentation output image here\n(bounding boxes on inscription image)",
             Inches(7.2), Inches(1.6), Inches(5.8), Inches(2.2),
             font_size=12, italic=True, color=GOLD, align=PP_ALIGN.CENTER)

divider(slide, 4.0)

add_text_box(slide, "Comparison with Related Work:",
             Inches(0.4), Inches(4.1), Inches(12.5), Inches(0.38),
             font_size=15, bold=True, color=DARK_BLUE)

add_table(slide,
    ["Study", "Script", "Method", "Classes", "Accuracy"],
    [
        ["Murugan & Visalakshi [1]",    "Ancient Tamil", "DRL Framework",    "—",   "—"],
        ["Gayathri Devi et al. [2]",    "Tamil Palm Leaf","CNN",             "—",   "—"],
        ["Vijaya Arjunan et al. [3]",   "Vatteluttu",    "Deep Learning",    "28",  "98.0%"],
        ["Howard et al. [25]",          "Generic",       "MobileNet",        "1000","70.6%"],
        ["VattalettuX (Proposed)",      "Vatteluttu",    "ResNet CNN [4]",   "247", "92.8%"],
    ],
    left=Inches(0.4), top=Inches(4.55), width=Inches(12.5), height=Inches(2.0))

add_text_box(slide,
    "Key Insight: Vijaya Arjunan achieved 98% on 28 classes. We achieve 92.8% on 247 classes — a 9× harder problem.\n"
    "Processing Speed: Avg. 1.8 sec/image on CPU  |  < 0.5 sec with GPU  |  Only fully deployed web app in comparison.",
    Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.75),
    font_size=12, italic=True, color=DARK_BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CONCLUSION & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "Conclusion and Future Work", "VattalettuX — Key Achievements & Roadmap")

add_text_box(slide, "✅  Key Achievements:",
             Inches(0.4), Inches(1.38), Inches(12.5), Inches(0.38),
             font_size=16, bold=True, color=DARK_BLUE)

achievements = [
    ("1", "Largest Vatteluttu character set",
     "247 classes — far more than the 28 classes covered by any prior study"),
    ("2", "Effective synthetic data pipeline",
     "2,47,000 images generated simulating real stone surface conditions"),
    ("3", "92.8% Top-1 accuracy",
     "ResNet CNN trained with PyTorch achieves strong performance across all 5 character categories"),
    ("4", "Fully deployed web application",
     "Accessible to historians and the general public — no installation required"),
]
y = 1.82
for num, title, detail in achievements:
    box = slide.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE if int(num)%2==1 else WHITE
    box.line.color.rgb = MID_BLUE
    add_text_box(slide, f"  {num}.  {title}  —  {detail}",
                 Inches(0.5), Inches(y+0.1), Inches(12.3), Inches(0.55),
                 font_size=13, bold=False, color=DARK_BLUE)
    y += 0.75

divider(slide, 4.95)

add_text_box(slide, "🚀  Future Work:",
             Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.38),
             font_size=16, bold=True, color=DARK_BLUE)

bullet_box(slide, [
    "▸  Transformer-based recognition (TrOCR) — for contextual word-level sequence recognition",
    "▸  GAN-based data synthesis — more realistic stone-texture training samples",
    "▸  Mobile deployment — MobileNet for field use by archaeologists",
    "▸  Real inscription dataset — label actual Vatteluttu stone photographs for better robustness",
    "▸  Language model post-correction — Tamil LM to auto-correct ambiguous OCR outputs",
],
    left=Inches(0.4), top=Inches(5.5), width=Inches(12.5), height=Inches(1.85),
    font_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – REFERENCES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
white_bg(slide)
slide_header(slide, "References", "")

refs = [
    "[1]  B. Murugan & P. Visalakshi — Ancient Tamil Inscription Recognition (DRL Framework), Heritage Science, 2024",
    "[2]  S. Gayathri Devi et al. — Deep Learning for Cursive Tamil Characters in Palm Leaf Manuscripts, Comput. Intelligence, 2022",
    "[3]  R. Vijaya Arjunan et al. — Deciphering Vatteluttu (Deep Learning, 28 classes, 98%), JISIS, 2025",
    "[4]  K. He et al. — Deep Residual Learning for Image Recognition (ResNet), CVPR, 2016",
    "[5]  K. Simonyan & A. Zisserman — Very Deep Convolutional Networks (VGGNet), ICLR, 2015",
    "[6]  A. Krizhevsky et al. — ImageNet Classification with Deep CNNs (AlexNet), NIPS, 2012",
    "[7]  Y. LeCun et al. — Gradient-Based Learning Applied to Document Recognition (LeNet), IEEE, 1998",
    "[8]  N. Otsu — A Threshold Selection Method from Gray-Level Histograms, IEEE Trans. SMC, 1979",
    "[9]  S. Suzuki & K. Abe — Topological Structural Analysis by Border Following (CCA), CVGIP, 1985",
    "[10] N. Srivastava et al. — Dropout: A Simple Way to Prevent Neural Networks from Overfitting, JMLR, 2014",
    "[11] D. P. Kingma & J. Ba — Adam: A Method for Stochastic Optimization, ICLR, 2015",
    "[12] S. Ioffe & C. Szegedy — Batch Normalization: Accelerating Deep Network Training, ICML, 2015",
    "[15] A. Paszke et al. — PyTorch: An Imperative Style Deep Learning Library, NeurIPS, 2019",
    "[16] C. Shorten & T. M. Khoshgoftaar — A Survey on Image Data Augmentation for Deep Learning, J. Big Data, 2019",
    "[17] G. Bradski — The OpenCV Library, Dr. Dobb's Journal, 2000",
    "[25] A. G. Howard et al. — MobileNets: Efficient CNNs for Mobile Vision Applications, arXiv, 2017",
]

txBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.38), Inches(12.5), Inches(5.9))
tf    = txBox.text_frame
tf.word_wrap = True

for i, ref in enumerate(refs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text           = ref
    run.font.name      = FONT
    run.font.size      = Pt(11.5)
    run.font.color.rgb = BLACK


# ─── SAVE ────────────────────────────────────────────────────────────────────
output_path = r"f:\final mca project\VattalettuX\VattalettuX_Presentation.pptx"
prs.save(output_path)
print(f"[DONE] Presentation saved: {output_path}")
print(f"       Total slides: {len(prs.slides)}")
